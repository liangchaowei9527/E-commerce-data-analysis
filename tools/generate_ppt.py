from __future__ import annotations

import json
from pathlib import Path

from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


PROJECT_ROOT = Path(__file__).resolve().parents[1]
PROCESSED_DIR = PROJECT_ROOT / "data" / "processed"
REPORTS_DIR = PROJECT_ROOT / "reports"


def add_title_slide(prs: Presentation, title: str, subtitle: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle


def add_bullet_slide(prs: Presentation, title: str, bullets: list[str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    text_frame = slide.placeholders[1].text_frame
    text_frame.clear()
    for idx, bullet in enumerate(bullets):
        paragraph = text_frame.paragraphs[0] if idx == 0 else text_frame.add_paragraph()
        paragraph.text = bullet
        paragraph.level = 0
        paragraph.font.size = Pt(20)


def add_image_slide(prs: Presentation, title: str, image_path: Path, caption: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title
    slide.shapes.add_picture(str(image_path), Inches(0.8), Inches(1.3), width=Inches(8.3))
    caption_box = slide.shapes.add_textbox(Inches(0.8), Inches(6.4), Inches(8.3), Inches(0.5))
    paragraph = caption_box.text_frame.paragraphs[0]
    paragraph.text = caption
    paragraph.font.size = Pt(18)
    paragraph.alignment = PP_ALIGN.LEFT


def main() -> None:
    summary = json.loads((PROCESSED_DIR / "analysis_summary.json").read_text(encoding="utf-8"))
    figures = {key: Path(value) for key, value in summary["figures"].items()}
    overview = summary["overview"]
    funnel = summary["funnel"]
    recommendations = summary["recommendations"]

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    add_title_slide(
        prs,
        "淘宝用户行为转化漏斗诊断与用户分层分析",
        "基于阿里云天池 UserBehavior 数据集",
    )
    add_bullet_slide(
        prs,
        "项目背景与目标",
        [
            "围绕浏览-收藏/加购-购买漏斗，定位关键转化损失环节。",
            "识别高活跃低转化用户群，为运营动作提供优先级。",
            "项目交付包括 Python 分析、SQL 验证、图表报告与面试汇报 PPT。",
        ],
    )
    add_bullet_slide(
        prs,
        "数据说明",
        [
            f"处理行为总量：{overview['processed_rows']:,} 条",
            f"覆盖独立用户：{overview['unique_users']:,} 人",
            f"行为结构：pv={overview['pv_events']:,}, fav={overview['fav_events']:,}, cart={overview['cart_events']:,}, buy={overview['buy_events']:,}",
            "清洗动作：去重、过滤无效行为、时间戳转北京时间、衍生时间字段。",
        ],
    )
    add_image_slide(prs, "整体行为概览", figures["behavior_distribution"], "先看流量大盘，确认行为结构以浏览为主。")
    add_image_slide(prs, "转化漏斗诊断", figures["funnel"], f"浏览转购买率 {funnel['pv_to_buy_rate']:.2%}，加购转购买率 {funnel['cart_to_buy_rate']:.2%}。")
    add_image_slide(prs, "分时段关键发现", figures["hourly_trend"], "对比高流量时段与高转化时段，识别流量承接问题。")
    add_image_slide(prs, "重点类目转化", figures["category_conversion"], "高曝光低成交类目是优化详情页和价格策略的优先对象。")
    add_image_slide(prs, "用户分层结果", figures["segment_distribution"], "高活跃低转化用户是最值得优先运营的人群。")
    add_bullet_slide(
        prs,
        "策略建议",
        [
            f"{item['priority']}：{item['issue']} - {item['action']}"
            for item in recommendations[:4]
        ],
    )
    add_bullet_slide(
        prs,
        "项目亮点与价值",
        [
            "使用分块计算处理大体量行为日志，避免一次性读入内存。",
            "SQL 与 Python 双栈交叉验证关键漏斗指标。",
            "输出从问题、证据到动作的完整业务分析闭环。",
        ],
    )
    add_bullet_slide(prs, "Q&A", ["谢谢观看", "可进一步扩展 Dashboard、活动分析与召回策略实验设计。"])

    REPORTS_DIR.mkdir(parents=True, exist_ok=True)
    prs.save(REPORTS_DIR / "电商用户行为转化诊断汇报PPT.pptx")


if __name__ == "__main__":
    main()
